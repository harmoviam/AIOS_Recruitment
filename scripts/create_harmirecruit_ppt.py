#!/usr/bin/env python3
"""Generate HarmiRecruit team-role pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette (teal + navy — WhatsApp-adjacent, not purple/cream cliché) ──
NAVY = RGBColor(0x0B, 0x1F, 0x33)
NAVY_MID = RGBColor(0x12, 0x2F, 0x4A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LIGHT = RGBColor(0x14, 0xB8, 0xA6)
CORAL = RGBColor(0xF0, 0x73, 0x4A)
GOLD = RGBColor(0xF5, 0xB8, 0x42)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF4, 0xF7, 0xF9)
SLATE = RGBColor(0x3D, 0x4F, 0x60)
MUTED = RGBColor(0x6B, 0x7C, 0x8C)
LIGHT_TEAL_BG = RGBColor(0xE6, 0xF6, 0xF4)

W, H = Inches(13.333), Inches(7.5)  # widescreen 16:9


def set_run(run, size=18, bold=False, color=NAVY, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text(shape, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT, font="Calibri"):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return tf


def add_para(tf, text, size=16, bold=False, color=SLATE, space_before=6, space_after=2, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill_solid(s, color)
    return s


def round_rect(slide, left, top, width, height, color, corner=0.1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    fill_solid(s, color)
    s.adjustments[0] = corner
    return s


def accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.55), color=TEAL):
    return rect(slide, left, top, width, height, color)


def star_pill(slide, left, top, label="KEY DIFFERENTIATOR"):
    pill = round_rect(slide, left, top, Inches(2.4), Inches(0.32), GOLD, corner=0.5)
    add_text(pill, f"★  {label}", size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    return pill


def section_header(slide, eyebrow, title, subtitle=None, dark=False):
    """Top-of-slide header block."""
    bg = NAVY if dark else OFF_WHITE
    fg = WHITE if dark else NAVY
    sub = RGBColor(0xA8, 0xC5, 0xC0) if dark else MUTED
    rect(slide, 0, 0, W, Inches(1.55) if subtitle else Inches(1.25), bg)
    rect(slide, 0, Inches(1.55) if subtitle else Inches(1.25), W, Inches(0.06), TEAL)

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(12), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = eyebrow.upper()
    set_run(run, size=11, bold=True, color=TEAL_LIGHT if dark else TEAL)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    run2 = p2.add_run()
    run2.text = title
    set_run(run2, size=28, bold=True, color=fg, font="Calibri")

    if subtitle:
        p3 = tf.add_paragraph()
        p3.space_before = Pt(4)
        run3 = p3.add_run()
        run3.text = subtitle
        set_run(run3, size=14, color=sub)


def bullet_card(slide, left, top, width, height, title, bullets, accent=TEAL, star=False):
    card = round_rect(slide, left, top, width, height, WHITE, corner=0.06)
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xED)
    card.line.width = Pt(1)
    accent_bar(slide, left + Inches(0.18), top + Inches(0.22), Inches(0.07), Inches(0.36), accent)

    title_box = slide.shapes.add_textbox(left + Inches(0.4), top + Inches(0.18), width - Inches(0.55), Inches(0.4))
    t = title
    if star:
        t = f"★  {title}"
    add_text(title_box, t, size=15, bold=True, color=NAVY)

    body = slide.shapes.add_textbox(left + Inches(0.28), top + Inches(0.6), width - Inches(0.5), height - Inches(0.75))
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        p.space_after = Pt(1)
        p.level = 0
        run = p.add_run()
        run.text = f"•  {b}"
        set_run(run, size=12, color=SLATE)
    return card


def footer(slide, page, total):
    rect(slide, 0, Inches(7.15), W, Inches(0.35), NAVY)
    left = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(6), Inches(0.25))
    add_text(left, "HarmiRecruit  ·  AI-first recruitment for India", size=10, color=RGBColor(0x8A, 0xA8, 0xB8))
    right = slide.shapes.add_textbox(Inches(10.5), Inches(7.2), Inches(2.3), Inches(0.25))
    add_text(right, f"{page}  /  {total}", size=10, color=RGBColor(0x8A, 0xA8, 0xB8), align=PP_ALIGN.RIGHT)


def new_slide(prs):
    blank = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    total = 14
    page = 0

    def next_page():
        nonlocal page
        page += 1
        return page

    # ═══════════════════════════════════════════════════════════
    # 1. TITLE
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    next_page()  # title = page 1 (no footer)
    rect(s, 0, 0, W, H, NAVY)
    # decorative side panel
    rect(s, 0, 0, Inches(0.22), H, TEAL)
    # soft accent band
    rect(s, 0, Inches(5.9), W, Inches(1.6), NAVY_MID)

    brand = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(1))
    add_text(brand, "HarmiRecruit", size=48, bold=True, color=WHITE, font="Calibri")

    tag = s.shapes.add_textbox(Inches(0.9), Inches(2.55), Inches(11.5), Inches(0.6))
    add_text(tag, "What Each Person on Your Team Gets", size=28, bold=False, color=TEAL_LIGHT)

    sub = s.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(10.5), Inches(1.2))
    tf = sub.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "AI-first recruitment with WhatsApp at the core — built for Indian agencies, staffing companies, and hiring teams."
    set_run(run, size=16, color=RGBColor(0xC5, 0xD4, 0xDE))

    chips = [
        ("👤  Recruiters", TEAL),
        ("👔  Hiring Managers", CORAL),
        ("🏢  Owners & Admins", GOLD),
        ("🙋  Candidates", TEAL_LIGHT),
    ]
    x = Inches(0.9)
    for label, color in chips:
        chip = round_rect(s, x, Inches(6.25), Inches(2.7), Inches(0.45), color, corner=0.4)
        add_text(chip, label, size=12, bold=True, color=NAVY if color == GOLD else WHITE, align=PP_ALIGN.CENTER)
        x += Inches(2.9)

    # ═══════════════════════════════════════════════════════════
    # 2. AGENDA / PROMISE
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(s, "The promise", "One platform. Four seats. Zero friction.", "Everyone on your team gets tools that match how they actually work.")

    promises = [
        ("👤  Recruiters", "Spend your day closing candidates,\nnot doing data entry", TEAL),
        ("👔  Hiring Managers", "See your whole team's desk\nat a glance", CORAL),
        ("🏢  Owners / Admins", "Run the business,\nnot the spreadsheet", GOLD),
        ("🙋  Candidates", "Zero friction —\nno app to install", TEAL_LIGHT),
    ]
    for i, (role, quote, color) in enumerate(promises):
        col = i % 2
        row = i // 2
        left = Inches(0.55) + col * Inches(6.35)
        top = Inches(2.0) + row * Inches(2.35)
        card = round_rect(s, left, top, Inches(6.05), Inches(2.1), WHITE, corner=0.06)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xED)
        rect(s, left, top, Inches(0.12), Inches(2.1), color)

        rb = s.shapes.add_textbox(left + Inches(0.4), top + Inches(0.35), Inches(5.3), Inches(0.4))
        add_text(rb, role, size=14, bold=True, color=MUTED)

        qb = s.shapes.add_textbox(left + Inches(0.4), top + Inches(0.85), Inches(5.3), Inches(1))
        tf = qb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f'"{quote}"'
        set_run(run, size=18, bold=True, color=NAVY)

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 3. RECRUITERS DIVIDER
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.22), H, TEAL)

    eye = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11), Inches(0.4))
    add_text(eye, "FOR RECRUITERS", size=14, bold=True, color=TEAL_LIGHT)

    title = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.2))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '"Spend your day closing candidates,\nnot doing data entry"'
    set_run(run, size=32, bold=True, color=WHITE)

    items = [
        "Candidate sourcing & intake",
        "Pipeline management",
        "WhatsApp built in ★",
        "Interviews",
        "Follow-up engine ★",
    ]
    box = s.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(11), Inches(1.8))
    tf = box.text_frame
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"→   {item}"
        set_run(run, size=16, color=RGBColor(0xC5, 0xD4, 0xDE))

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 4. SOURCING & INTAKE
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(s, "Recruiters · Sourcing & Intake", "Drop a resume. Get a full profile. Done.", "From PDF to searchable candidate in seconds — no typing.")

    bullets_left = [
        "Drop a resume (PDF/Word) → AI extracts contact, skills, experience, education, salary, notice period",
        "Applications flow in automatically from your branded careers page — no manual entry",
        "Duplicate detection: same candidate applying twice is merged, never duplicated",
    ]
    bullets_right = [
        "Bulk CSV import with validation for migrating your existing database",
        "Search across everything — name, phone, skills, and even inside resume text",
        '"React Mohali" finds them — full-text search that actually works',
    ]
    bullet_card(s, Inches(0.5), Inches(1.9), Inches(6.05), Inches(4.7), "Smart intake", bullets_left, TEAL)
    bullet_card(s, Inches(6.8), Inches(1.9), Inches(6.05), Inches(4.7), "Search & migrate", bullets_right, CORAL)
    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 5. PIPELINE
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(s, "Recruiters · Pipeline", "Your desk, visualised — drag, drop, close.", "One candidate. Multiple jobs. Each with its own stage.")

    # stage flow
    stages = ["Applied", "Screening", "Interview", "Selected", "Joined"]
    sx = Inches(0.55)
    for i, stg in enumerate(stages):
        chip = round_rect(s, sx, Inches(1.9), Inches(2.15), Inches(0.5), TEAL if i < 4 else CORAL, corner=0.35)
        add_text(chip, stg, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            arrow = s.shapes.add_textbox(sx + Inches(2.05), Inches(1.95), Inches(0.35), Inches(0.4))
            add_text(arrow, "→", size=16, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        sx += Inches(2.5)

    feats = [
        ("Multi-job submit", "Submit one candidate to multiple jobs at once — each submission tracks its own stage"),
        ("AI Match Score /10", "Every candidate ranked against the job, with strengths and gaps explained"),
        ("Hot + Bulk actions", "Flag hot candidates, bulk stage updates, one-click CSV export"),
        ("360° profile", "Timeline, conversations, interviews, scores, notes — all on one screen"),
    ]
    for i, (t, d) in enumerate(feats):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * Inches(6.4)
        top = Inches(2.75) + row * Inches(1.9)
        card = round_rect(s, left, top, Inches(6.15), Inches(1.7), WHITE, corner=0.06)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xED)
        accent_bar(s, left + Inches(0.2), top + Inches(0.35), Inches(0.07), Inches(0.9), TEAL if i % 2 == 0 else CORAL)
        tb = s.shapes.add_textbox(left + Inches(0.45), top + Inches(0.3), Inches(5.4), Inches(0.4))
        add_text(tb, t, size=16, bold=True, color=NAVY)
        db = s.shapes.add_textbox(left + Inches(0.45), top + Inches(0.8), Inches(5.4), Inches(0.7))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = d
        set_run(run, size=13, color=SLATE)

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 6. WHATSAPP
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(
        s,
        "Recruiters · WhatsApp ★",
        "Inbox inside the platform. Candidates where they already are.",
        "Official Meta Business API — send, receive, and follow up without leaving HarmiRecruit.",
    )
    star_pill(s, Inches(10.5), Inches(0.35), "STAR FEATURE")

    wa_feats = [
        ("Full WhatsApp inbox", "Send and receive on the official Meta Business API — your conversations live inside the ATS"),
        ("AI reply suggestions", "Drafted responses for every conversation — edit and send in one tap"),
        ("Interview & join links", "Invites, join links, and follow-ups delivered straight to the candidate's WhatsApp"),
        ("Recruiter signatures", "Workspace branding on every message — your name, your agency, your professionalism"),
    ]
    for i, (t, d) in enumerate(wa_feats):
        top = Inches(1.9) + i * Inches(1.2)
        card = round_rect(s, Inches(0.5), top, Inches(12.3), Inches(1.05), WHITE, corner=0.05)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xED)
        num = round_rect(s, Inches(0.7), top + Inches(0.25), Inches(0.55), Inches(0.55), TEAL, corner=0.3)
        add_text(num, str(i + 1), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb = s.shapes.add_textbox(Inches(1.5), top + Inches(0.18), Inches(10.8), Inches(0.35))
        add_text(tb, t, size=16, bold=True, color=NAVY)
        db = s.shapes.add_textbox(Inches(1.5), top + Inches(0.55), Inches(10.8), Inches(0.4))
        add_text(db, d, size=13, color=SLATE)

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 7. INTERVIEWS
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(s, "Recruiters · Interviews", "Two clicks to schedule. Zero Zoom licences.", "WhatsApp + email + calendar invite — candidate ready, you ready.")

    interview_cards = [
        ("Schedule fast", [
            "Two-click scheduling",
            "WhatsApp + email with calendar invite (.ics)",
            "Candidate gets everything instantly",
        ]),
        ("Built-in video room", [
            "No Zoom licence needed",
            "Mic check before joining",
            "One-tap candidate join link in browser",
        ]),
        ("First-call scorecard", [
            "Rate energy, motivation, communication",
            "Red-flag signals captured",
            "Automatic risk rating",
        ]),
        ("AI screening Qs", [
            "Questions tailored to each job's JD",
            "Ready before you dial",
            "Consistent quality across the team",
        ]),
    ]
    for i, (t, bullets) in enumerate(interview_cards):
        left = Inches(0.4) + i * Inches(3.2)
        bullet_card(s, left, Inches(1.9), Inches(3.05), Inches(4.7), t, bullets, [TEAL, CORAL, GOLD, TEAL_LIGHT][i])

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 8. FOLLOW-UP ENGINE
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(
        s,
        "Recruiters · Follow-up Engine ★",
        "The dropout killer.",
        "From offer to Day-90 — automatic schedules so candidates don't ghost.",
    )
    star_pill(s, Inches(10.5), Inches(0.35), "STAR FEATURE")

    milestones = [
        ("Offer", "Selected → offer stage kicks off"),
        ("Join −N", "Joining-day countdown"),
        ("Day 1", "Onboarding check-in"),
        ("Day 30", "Settling-in pulse"),
        ("Day 60", "Stay-risk check"),
        ("Day 90", "Retention milestone"),
    ]
    for i, (m, d) in enumerate(milestones):
        left = Inches(0.4) + i * Inches(2.15)
        card = round_rect(s, left, Inches(1.95), Inches(2.0), Inches(1.55), WHITE, corner=0.08)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xED)
        top_bar = rect(s, left, Inches(1.95), Inches(2.0), Inches(0.12), TEAL if i < 3 else CORAL)
        mb = s.shapes.add_textbox(left + Inches(0.1), Inches(2.25), Inches(1.8), Inches(0.4))
        add_text(mb, m, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        db = s.shapes.add_textbox(left + Inches(0.1), Inches(2.7), Inches(1.8), Inches(0.6))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = d
        set_run(run, size=11, color=SLATE)

    extras = [
        ("AI scripts & messages", "Call scripts and WhatsApp messages generated for every follow-up milestone"),
        ("Escalation alerts", "Know before they ghost — escalations when a candidate stops responding"),
        ("Set and forget", "Schedule starts the moment a candidate is selected — your team just executes"),
    ]
    for i, (t, d) in enumerate(extras):
        left = Inches(0.4) + i * Inches(4.25)
        card = round_rect(s, left, Inches(3.85), Inches(4.05), Inches(2.7), NAVY if i == 1 else WHITE, corner=0.06)
        if i != 1:
            card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xED)
        tb = s.shapes.add_textbox(left + Inches(0.3), Inches(4.15), Inches(3.45), Inches(0.5))
        add_text(tb, t, size=15, bold=True, color=WHITE if i == 1 else NAVY)
        db = s.shapes.add_textbox(left + Inches(0.3), Inches(4.75), Inches(3.45), Inches(1.4))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = d
        set_run(run, size=13, color=RGBColor(0xC5, 0xD4, 0xDE) if i == 1 else SLATE)

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 9. HIRING MANAGERS
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(
        s,
        "For Hiring Managers / Team Leads",
        '"See your whole team\'s desk at a glance"',
        "Lead the desk. Spot bottlenecks. Ship better shortlists.",
    )

    hm = [
        ("HM Dashboard", "Team pipeline, per-recruiter workload, and today's interviews — one screen"),
        ("Build your team", "Add recruiters; candidates auto-scoped to your team"),
        ("Switch views", "My candidates ↔ My team's candidates — toggle instantly"),
        ("AI JD Generator", "Create & edit jobs with role-specific templates for IT, BPO, banking & more"),
        ("Review before client", "Interview evaluations & screening scorecards before shortlisting"),
        ("Team performance", "Submissions, interviews, selects, joins — per recruiter reports"),
    ]
    for i, (t, d) in enumerate(hm):
        col = i % 3
        row = i // 3
        left = Inches(0.4) + col * Inches(4.25)
        top = Inches(1.9) + row * Inches(2.4)
        card = round_rect(s, left, top, Inches(4.05), Inches(2.15), WHITE, corner=0.06)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xED)
        accent_bar(s, left + Inches(0.2), top + Inches(0.3), Inches(0.07), Inches(0.45), CORAL)
        tb = s.shapes.add_textbox(left + Inches(0.4), top + Inches(0.28), Inches(3.4), Inches(0.4))
        add_text(tb, t, size=15, bold=True, color=NAVY)
        db = s.shapes.add_textbox(left + Inches(0.25), top + Inches(0.85), Inches(3.55), Inches(1.1))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = d
        set_run(run, size=13, color=SLATE)

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 10. ADMIN / OWNER DIVIDER + CONTENT
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(
        s,
        "For Agency Owner / Admin",
        '"Run the business, not the spreadsheet"',
        "KPIs, clients, branding, billing — the control room for your agency.",
    )

    admin_left = [
        "Admin dashboard: org-wide KPIs, funnel, recruiter leaderboard, live activity feed",
        "Client (company) management with locations + geo-matching for nearby candidates",
        "Team management: recruiters, HMs, admins — automatic welcome email for every new user",
        "Branded careers page at your own link — copy apply links to WhatsApp, LinkedIn, Naukri",
    ]
    admin_right = [
        "Reports centre: productivity, funnel, offer-to-join — export to CSV/Excel",
        "Full analytics: source mix, stage conversion, trends over time",
        "Workspace branding: logo, colours, recruiter WhatsApp signatures",
        "Billing self-service via Razorpay (UPI, cards, netbanking) + GST invoices",
        "Data isolation: your workspace is completely private to your organisation",
    ]
    bullet_card(s, Inches(0.45), Inches(1.9), Inches(6.15), Inches(4.7), "Operate & grow", admin_left, GOLD)
    bullet_card(s, Inches(6.8), Inches(1.9), Inches(6.05), Inches(4.7), "Measure & monetise", admin_right, TEAL)
    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 11. CANDIDATES
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(
        s,
        "For Candidates",
        "Zero friction. No app to install.",
        "Apply from any phone. Interview from the browser. That's it.",
    )

    steps = [
        ("01", "Apply", "Branded careers page — name, number, resume, done"),
        ("02", "Confirm", "Instant application-received email"),
        ("03", "Invite", "WhatsApp + email with calendar file on their phone"),
        ("04", "Join", "One-tap video interview — browser only, mic check included"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        left = Inches(0.45) + i * Inches(3.2)
        card = round_rect(s, left, Inches(2.2), Inches(3.0), Inches(4.2), WHITE, corner=0.08)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xED)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.95), Inches(2.6), Inches(1.1), Inches(1.1))
        fill_solid(circle, TEAL if i < 3 else CORAL)
        add_text(circle, num, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # fix vertical centering for oval text
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle.text_frame.word_wrap = False

        tb = s.shapes.add_textbox(left + Inches(0.2), Inches(4.0), Inches(2.6), Inches(0.5))
        add_text(tb, title, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        db = s.shapes.add_textbox(left + Inches(0.25), Inches(4.6), Inches(2.5), Inches(1.4))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        set_run(run, size=13, color=SLATE)

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 12. AI LAYER
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(
        s,
        "The AI Layer",
        "Works across every role — quietly doing the heavy lifting.",
        "Six AI capabilities. One outcome: faster, sharper hiring.",
    )

    ai_rows = [
        ("Resume Parser", "Resume → complete profile in seconds"),
        ("Match Score", "Ranks every candidate against the job, with reasons"),
        ("JD Generator", "Complete, polished job descriptions from just a title"),
        ("Screening Q Generator", "Interview questions derived from each JD"),
        ("WhatsApp Reply Suggestions", "AI-drafted responses in every conversation"),
        ("Follow-up Script Writer", "Call scripts & messages for every follow-up milestone"),
    ]

    # table header
    hdr = round_rect(s, Inches(0.55), Inches(1.85), Inches(12.2), Inches(0.55), NAVY, corner=0.04)
    hb = s.shapes.add_textbox(Inches(0.85), Inches(1.95), Inches(4.5), Inches(0.4))
    add_text(hb, "AI Feature", size=14, bold=True, color=WHITE)
    hb2 = s.shapes.add_textbox(Inches(5.5), Inches(1.95), Inches(7), Inches(0.4))
    add_text(hb2, "What it does for you", size=14, bold=True, color=WHITE)

    for i, (feat, does) in enumerate(ai_rows):
        top = Inches(2.5) + i * Inches(0.7)
        bg = WHITE if i % 2 == 0 else LIGHT_TEAL_BG
        row = rect(s, Inches(0.55), top, Inches(12.2), Inches(0.7), bg)
        fb = s.shapes.add_textbox(Inches(0.85), top + Inches(0.15), Inches(4.5), Inches(0.4))
        add_text(fb, feat, size=14, bold=True, color=NAVY)
        db = s.shapes.add_textbox(Inches(5.5), top + Inches(0.15), Inches(7), Inches(0.4))
        add_text(db, does, size=14, color=SLATE)

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 13. WHY IT CLICKS (summary)
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, OFF_WHITE)
    section_header(s, "Why teams switch", "Built for how Indian recruiting actually works.", "WhatsApp-native. AI-assisted. Agency-ready.")

    reasons = [
        ("WhatsApp at the core", "Candidates already live there — so does your outreach, invites, and follow-ups"),
        ("AI that saves hours", "Parsing, matching, JDs, scripts, and replies — less admin, more closing"),
        ("Role-perfect seats", "Recruiters, HMs, and owners each get a dashboard that fits their job"),
        ("Careers → Pipeline", "Branded apply page feeds straight into your ATS with dedupe"),
        ("Join, don't ghost", "Offer-to-join follow-ups with escalations before dropouts happen"),
        ("Pay the Indian way", "Razorpay: UPI, cards, netbanking, GST invoices, monthly or annual"),
    ]
    for i, (t, d) in enumerate(reasons):
        col = i % 3
        row = i // 3
        left = Inches(0.4) + col * Inches(4.25)
        top = Inches(1.9) + row * Inches(2.4)
        card = round_rect(s, left, top, Inches(4.05), Inches(2.15), WHITE, corner=0.06)
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xED)
        n = round_rect(s, left + Inches(0.25), top + Inches(0.25), Inches(0.45), Inches(0.45), TEAL, corner=0.3)
        add_text(n, str(i + 1), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb = s.shapes.add_textbox(left + Inches(0.85), top + Inches(0.3), Inches(2.9), Inches(0.4))
        add_text(tb, t, size=14, bold=True, color=NAVY)
        db = s.shapes.add_textbox(left + Inches(0.25), top + Inches(0.95), Inches(3.55), Inches(1.0))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = d
        set_run(run, size=13, color=SLATE)

    footer(s, pnum, total)

    # ═══════════════════════════════════════════════════════════
    # 14. CLOSING
    # ═══════════════════════════════════════════════════════════
    s = new_slide(prs)
    pnum = next_page()
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.22), H, TEAL)
    rect(s, 0, Inches(5.8), W, Inches(1.7), NAVY_MID)

    eye = s.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11), Inches(0.4))
    add_text(eye, "READY WHEN YOU ARE", size=13, bold=True, color=TEAL_LIGHT)

    title = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.5))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Give every seat on your team\nthe tools to hire faster."
    set_run(run, size=34, bold=True, color=WHITE)

    sub = s.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(11), Inches(0.6))
    add_text(sub, "HarmiRecruit — AI-first recruitment with WhatsApp at the core.", size=16, color=RGBColor(0xC5, 0xD4, 0xDE))

    cta = round_rect(s, Inches(0.9), Inches(6.2), Inches(3.2), Inches(0.55), TEAL, corner=0.4)
    add_text(cta, "Book a demo →", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    side = s.shapes.add_textbox(Inches(4.5), Inches(6.3), Inches(7.5), Inches(0.4))
    add_text(side, "Recruiters  ·  Hiring Managers  ·  Owners  ·  Candidates", size=13, color=RGBColor(0x8A, 0xA8, 0xB8))

    out = "/Users/jyotiranjan/workarea/projects/AIOS_Recruitment/docs/HarmiRecruit_What_Each_Person_Gets.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    build()
